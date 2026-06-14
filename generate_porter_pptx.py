
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_porter_pptx(filename="porter_diagram.pptx"):
    prs = Presentation()
    # Set slide size to 16:9
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    slide_layout = prs.slide_layouts[6] # Blank slide
    slide = prs.slides.add_slide(slide_layout)

    # Colors
    dark_blue = RGBColor(44, 95, 141)
    medium_blue = RGBColor(74, 139, 191)
    light_blue = RGBColor(124, 185, 232)
    green = RGBColor(92, 184, 92)
    teal = RGBColor(76, 175, 80)
    white = RGBColor(255, 255, 255)
    black = RGBColor(51, 51, 51)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    tf.text = "الضغط التنافسي: قوى بورتر الخمسة"
    p = tf.paragraphs[0]
    p.font.bold = True
    p.font.size = Pt(32)
    p.alignment = PP_ALIGN.CENTER

    # Center circle
    center_x, center_y = Inches(6.66), Inches(3.75)
    circle_radius = Inches(1.2)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, center_x - circle_radius/2, center_y - circle_radius/2, circle_radius, circle_radius)
    circle.fill.solid()
    circle.fill.fore_color.rgb = dark_blue
    circle.line.color.rgb = white
    circle.line.width = Pt(2)
    
    # Center text
    text_box = slide.shapes.add_textbox(center_x - circle_radius/2, center_y - Inches(0.3), circle_radius, Inches(0.6))
    tf = text_box.text_frame
    tf.text = "مستشفى\nالشاعر الخاص"
    for p in tf.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        p.font.color.rgb = white
        p.font.size = Pt(16)
        p.font.bold = True

    # Helper function to add info box
    def add_info_box(x, y, header_text, body_text, color):
        width, height = Inches(3.5), Inches(1.8)
        header_height = Inches(0.5)
        
        # Header
        header = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, width, header_height)
        header.fill.solid()
        header.fill.fore_color.rgb = color
        header.line.width = Pt(0)
        
        h_tf = header.text_frame
        h_tf.text = header_text
        h_p = h_tf.paragraphs[0]
        h_p.alignment = PP_ALIGN.CENTER
        h_p.font.color.rgb = white
        h_p.font.size = Pt(14)
        h_p.font.bold = True
        
        # Body
        body = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y + header_height - Inches(0.1), width, height - header_height + Inches(0.1))
        body.fill.solid()
        body.fill.fore_color.rgb = white
        body.line.color.rgb = color
        body.line.width = Pt(1.5)
        
        b_tf = body.text_frame
        b_tf.word_wrap = True
        b_tf.text = body_text
        for p in b_tf.paragraphs:
            p.alignment = PP_ALIGN.RIGHT
            p.font.color.rgb = black
            p.font.size = Pt(12)

        return x + width/2, y + height/2

    # Box positions
    # Top Left
    tl_pos = add_info_box(Inches(0.5), Inches(1.2), "حدة المنافسة (3.5/5)", "منافسة قوية بين 3 منشآت خاصة و6 حكومية حول السعر والموقع. (الرقمنة هي المهرب الوحيد).", dark_blue)
    # Top Right
    tr_pos = add_info_box(Inches(9.3), Inches(1.2), "قوة المشترين (4/5)", "وعي المريض يرتفع، وحساسيته للجودة تتفوق على السعر في حال توفر التجربة الممتازة.", medium_blue)
    # Middle Right
    mr_pos = add_info_box(Inches(9.3), Inches(3.3), "قوة الموردين (3/5)", "قوة تفاوضية عالية للكوادر الطبية المتخصصة.", light_blue)
    # Bottom Right
    br_pos = add_info_box(Inches(9.3), Inches(5.4), "الداخلون الجدد (3/5)", "حواجز دخول مالية عالية (30 مليون+)، لكن التهديد قادم من العيادات الرقمية.", green)
    # Bottom Left
    bl_pos = add_info_box(Inches(0.5), Inches(5.4), "المنتجات البديلة (2.5/5)", "الطب التقليدي، والسفر للعلاج (المريض يفضل الرعاية المحلية إذا توفرت الجودة).", teal)

    # Add connectors (lines)
    def add_connector(start_x, start_y, end_x, end_y, color):
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, start_x, start_y, end_x, end_y)
        line.line.color.rgb = color
        line.line.width = Pt(3)

    # Connect center to boxes
    add_connector(center_x - Inches(0.6), center_y - Inches(0.3), tl_pos[0] + Inches(1), tl_pos[1], dark_blue)
    add_connector(center_x + Inches(0.6), center_y - Inches(0.3), tr_pos[0] - Inches(1), tr_pos[1], medium_blue)
    add_connector(center_x + Inches(0.6), center_y, mr_pos[0] - Inches(1), mr_pos[1], light_blue)
    add_connector(center_x + Inches(0.6), center_y + Inches(0.3), br_pos[0] - Inches(1), br_pos[1], green)
    add_connector(center_x - Inches(0.6), center_y + Inches(0.3), bl_pos[0] + Inches(1), bl_pos[1], teal)

    prs.save(filename)

if __name__ == "__main__":
    create_porter_pptx()
