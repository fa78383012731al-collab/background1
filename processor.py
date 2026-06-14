"""
processor.py — Parse PPTX, detect diagrams, orchestrate reconstruction.

Strategy:
  - Slide has SmartArt  → convert to editable shapes
  - Slide has native shape groups (many rects/connectors) → already editable, enrich
  - Slide has embedded picture diagrams → OpenCV reconstruct
  - Always preserve original content; add editable layer on top
"""
from __future__ import annotations

import io
import os
from pathlib import Path
from typing import Callable

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image

from rebuild_diagram import rebuild_diagram_on_slide


# ── Diagram detection ────────────────────────────────────────────────────────

def _has_smartart(slide) -> bool:
    """Check if slide contains SmartArt (graphicFrame with dgm namespace)."""
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.FREEFORM:
            return True
        el = shape._element
        tag = el.tag.split("}")[-1] if "}" in el.tag else el.tag
        if tag == "graphicFrame":
            xml = el.xml if hasattr(el, "xml") else str(el)
            if "dgm:" in xml or "drawingml/2006/diagram" in xml:
                return True
    return False


def _count_native_diagram_shapes(slide) -> int:
    """Count rectangle/connector shapes — high count = native diagram."""
    count = 0
    for shape in slide.shapes:
        if shape.shape_type in (
            MSO_SHAPE_TYPE.AUTO_SHAPE,
            MSO_SHAPE_TYPE.FREEFORM,
            MSO_SHAPE_TYPE.LINE,
            MSO_SHAPE_TYPE.TEXT_BOX,
        ):
            count += 1
    return count


def _is_diagram_image(img: Image.Image) -> tuple[bool, str]:
    """
    Broad heuristic for image-embedded diagrams.
    Diagrams: moderate non-white, limited unique colours.
    """
    img_rgb = img.convert("RGB")
    w, h = img_rgb.size
    if w < 80 or h < 80:
        return False, "too small"

    pixels = list(img_rgb.getdata())
    total = len(pixels)
    non_white = sum(1 for r, g, b in pixels if not (r > 220 and g > 220 and b > 220))
    non_white_ratio = non_white / total

    sampled = pixels[::max(1, total // 3000)]
    unique_colours = len(set(sampled))

    if non_white_ratio < 0.02:
        return False, "mostly blank"
    if non_white_ratio > 0.92 and unique_colours > 5000:
        return False, "likely photo"
    return True, f"diagram ({non_white_ratio:.0%} content, {unique_colours} colours)"


def _classify_diagram_type(img: Image.Image) -> str:
    w, h = img.size
    ratio = w / h if h else 1
    pixels = list(img.convert("RGB").getdata())[::max(1, (w * h) // 500)]
    unique = len(set(pixels))
    if unique < 20:
        return "Flowchart"
    if ratio > 2.0:
        return "Process / Timeline"
    if ratio < 0.5:
        return "Vertical Diagram"
    if unique < 80:
        return "Infographic"
    return "Diagram"


def _extract_picture_shapes(slide) -> list[dict]:
    images = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                blob = shape.image.blob
                img = Image.open(io.BytesIO(blob))
                images.append({
                    "shape": shape,
                    "image": img,
                    "left": shape.left,
                    "top": shape.top,
                    "width": shape.width,
                    "height": shape.height,
                })
            except Exception:
                pass
    return images


# ── Main processor ────────────────────────────────────────────────────────────

def process_pptx(
    pptx_path: str,
    output_dir: str,
    log: Callable[[str, int], None] | None = None,
) -> dict:
    def _log(msg: str, pct: int | None = None):
        print(msg)
        if log:
            log(msg, pct)

    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)

    _log("Opening PPTX…", 5)
    prs = Presentation(pptx_path)
    slides = prs.slides
    n_slides = len(slides)
    _log(f"Loaded {n_slides} slide(s). Slide size: {prs.slide_width} × {prs.slide_height} EMU", 8)

    slide_results = []
    all_diagram_count = 0

    for idx, slide in enumerate(slides):
        pct_start = 10 + int(65 * idx / max(n_slides, 1))
        _log(f"── Slide {idx+1}/{n_slides} ──", pct_start)

        diagrams_on_slide = []

        # ── A: Check for SmartArt ─────────────────────────────────────────
        has_sa = _has_smartart(slide)
        if has_sa:
            _log(f"  SmartArt detected on slide {idx+1}", pct_start + 2)

        # ── B: Check native shape groups ─────────────────────────────────
        native_count = _count_native_diagram_shapes(slide)
        _log(f"  Native diagram shapes: {native_count}", pct_start + 3)

        # ── C: Check for embedded picture diagrams ────────────────────────
        pics = _extract_picture_shapes(slide)
        _log(f"  Embedded pictures: {len(pics)}", pct_start + 4)

        if native_count >= 3 and not pics:
            # Already editable native shapes — document them
            _log(f"  Slide {idx+1}: {native_count} native shapes (already editable) — enriching…")
            diag_info = _enrich_native_shapes(slide, prs, idx, str(output_dir))
            if diag_info:
                diag_info["type"] = "SmartArt / Native Shapes"
                diag_info["description"] = f"{native_count} editable shapes"
                diagrams_on_slide.append(diag_info)
                all_diagram_count += 1

        for img_data in pics:
            img = img_data["image"]
            is_diag, reason = _is_diagram_image(img)
            _log(f"  Picture: {reason}")
            if not is_diag:
                continue

            diag_type = _classify_diagram_type(img)
            _log(f"  → Reconstructing as {diag_type}…")

            slide_w_emu = prs.slide_width
            slide_h_emu = prs.slide_height

            diag_info = rebuild_diagram_on_slide(
                slide=slide,
                img_data=img_data,
                diagram_type=diag_type,
                diagram_index=all_diagram_count,
                output_dir=str(output_dir),
                slide_w_emu=slide_w_emu,
                slide_h_emu=slide_h_emu,
            )
            diag_info["type"] = diag_type
            diag_info["description"] = reason
            diagrams_on_slide.append(diag_info)
            all_diagram_count += 1

        slide_results.append({"index": idx, "diagrams": diagrams_on_slide})

    _log(f"Total diagrams reconstructed: {all_diagram_count}", 82)

    pptx_out = str(output_dir / "reconstructed.pptx")
    _log("Saving PPTX…", 84)
    prs.save(pptx_out)

    # Export SVG/PNG from first preview
    svg_out = png_out = None
    from export_svg import export_slide_to_svg
    from export_png import export_slide_to_png
    first_diag = next((d for s in slide_results for d in s["diagrams"]), None)
    if first_diag and first_diag.get("preview_png"):
        try:
            _log("Exporting SVG…", 88)
            svg_out = export_slide_to_svg(first_diag, str(output_dir))
        except Exception as e:
            _log(f"SVG export skipped: {e}")
        try:
            _log("Exporting PNG…", 92)
            png_out = export_slide_to_png(first_diag, str(output_dir))
        except Exception as e:
            _log(f"PNG export skipped: {e}")

    _log("Done ✅", 98)
    return {
        "slides": slide_results,
        "pptx_out": pptx_out,
        "svg_out": svg_out,
        "png_out": png_out,
        "diagram_count": all_diagram_count,
    }


def _enrich_native_shapes(slide, prs, slide_idx: int, output_dir: str) -> dict | None:
    """Document native shapes and generate a preview PNG."""
    try:
        from PIL import Image, ImageDraw
        slide_w = prs.slide_width
        slide_h = prs.slide_height
        scale = 1200 / slide_w
        pw = int(slide_w * scale)
        ph = int(slide_h * scale)
        canvas = Image.new("RGB", (pw, ph), (255, 255, 255))
        draw = ImageDraw.Draw(canvas)

        from pptx.util import Emu
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        from pptx.dml.color import RGBColor

        for shape in slide.shapes:
            try:
                l = int(shape.left * scale) if shape.left else 0
                t = int(shape.top  * scale) if shape.top  else 0
                w = int(shape.width * scale) if shape.width else 0
                h = int(shape.height * scale) if shape.height else 0

                fill_color = (173, 216, 230)
                try:
                    if shape.fill.type is not None:
                        fc = shape.fill.fore_color.rgb
                        fill_color = (fc[0], fc[1], fc[2])
                except Exception:
                    pass

                if shape.shape_type in (MSO_SHAPE_TYPE.AUTO_SHAPE,
                                         MSO_SHAPE_TYPE.TEXT_BOX):
                    draw.rectangle([l, t, l+w, t+h],
                                   fill=(*fill_color, 180), outline=(80,80,80), width=2)

                    if shape.has_text_frame:
                        txt = shape.text_frame.text[:50]
                        if txt.strip():
                            draw.text((l+4, t+4), txt, fill=(20, 20, 20))
            except Exception:
                pass

        out_path = str(Path(output_dir) / f"preview_native_{slide_idx}.png")
        canvas.save(out_path)
        return {"preview_png": out_path}
    except Exception as e:
        print(f"  _enrich_native_shapes error: {e}")
        return None
